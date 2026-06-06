"""
MouthMap — Digital Twin Patient Portal
Presentation for Dr. Christian Brenes
Meeting: June 4, 2026 @ 4:00 PM (Zoom)

V7 — TIGHT 10-SLIDE VERSION. ~10 min presenting, 35 min conversation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from zipfile import ZipFile
import shutil, os

# ─── PATHS ───
TEMPLATE_SRC = r'C:\Users\Chris\Desktop\WEBSITES\THINK! VENTURES\docs\Think_Ventures_Template.potx'
TEMPLATE      = r'C:\Users\Chris\Desktop\WEBSITES\THINK! VENTURES\docs\Think_Ventures_Template_converted.pptx'
OUTPUT        = r'C:\Users\Chris\Downloads\MouthMap_DrBrenes_Presentation.pptx'
MM_LOGO       = r'C:\Users\Chris\.gemini\antigravity\brain\ce94a7df-3f74-4928-946a-59b87eb9ce19\mouthmap_logo_icon_1780578615383.png'

# ─── BRAND COLORS ───
TEAL       = RGBColor(0x0D, 0x4F, 0x4F)
GOLD       = RGBColor(0xF5, 0xA6, 0x23)
EMERALD    = RGBColor(0x10, 0xB9, 0x81)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x1A, 0x2A, 0x3A)
MUTED_TEXT = RGBColor(0x4A, 0x5A, 0x6A)
CARD_BG    = RGBColor(0xE8, 0xED, 0xF2)
LJ_BLUE    = RGBColor(0x4A, 0x7F, 0xB5)

# ─── CONVERT .potx → .pptx ───
if not os.path.exists(TEMPLATE):
    shutil.copy2(TEMPLATE_SRC, TEMPLATE)
    with ZipFile(TEMPLATE, 'r') as zin:
        all_data = {n: zin.read(n) for n in zin.namelist()}
    ct = all_data['[Content_Types].xml'].decode('utf-8').replace(
        'presentationml.template.main+xml',
        'presentationml.presentation.main+xml')
    all_data['[Content_Types].xml'] = ct.encode('utf-8')
    with ZipFile(TEMPLATE, 'w') as zout:
        for n, d in all_data.items():
            zout.writestr(n, d)

prs = Presentation(TEMPLATE)
LO_TITLE   = prs.slide_layouts[0]
LO_SECTION = prs.slide_layouts[2]

# ─── HELPERS ───

def txt(slide, left, top, width, height, text, size=18, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font_name=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
    p.font.color.rgb = color; p.font.bold = bold
    p.font.name = font_name if font_name else ("Outfit" if bold and size >= 20 else "Inter")
    p.alignment = align
    return box

def bullets(slide, left, top, width, height, items, size=15, color=DARK_TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.color.rgb = color
        p.font.name = "Inter"; p.space_after = Pt(6)
    return box

def card(slide, left, top, width, height, color=CARD_BG):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def teal_card(s, l, t, w, h): return card(s, l, t, w, h, TEAL)

def notes(slide, text):
    ns = slide.notes_slide; tf = ns.notes_text_frame
    if tf: tf.text = text

def add_logo(slide, left=10.5, top=1.3, height=1.2):
    if os.path.exists(MM_LOGO):
        slide.shapes.add_picture(MM_LOGO, Inches(left), Inches(top), height=Inches(height))


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE (30 sec)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_TITLE)

for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "MouthMap"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(52); p.font.bold = True; p.font.name = "Outfit"

txt(slide, 0.8, 3.4, 8, 0.6, "The Open-Architecture Patient Portal", 24, MUTED_TEXT)
txt(slide, 0.8, 4.1, 8, 0.4, "Navigate Your Oral Health.", 20, TEAL, True, font_name="Outfit")
txt(slide, 0.8, 4.8, 8, 0.4, "Prepared for Dr. Christian Brenes  |  June 4, 2026", 14, MUTED_TEXT)
txt(slide, 0.8, 5.3, 8, 0.4, "W. Chris Harrison  |  Think! Design and Planning, LLC", 14, GOLD, True)

if os.path.exists(MM_LOGO):
    slide.shapes.add_picture(MM_LOGO, Inches(9.0), Inches(2.5), height=Inches(2.5))

notes(slide, "Thank you for taking the time. MouthMap is a digital twin patient portal — a secure website where patients interact with their 3D scans and AI findings from home. This is an exploration, not a pitch. I need your expertise to tell me if this makes sense.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: THE WEBSITE — FREE (1 min)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_SECTION)

txt(slide, 0.8, 1.2, 5, 0.3, "A GIFT", 14, GOLD, True)
txt(slide, 0.8, 1.55, 10, 0.5, "The Website Is Yours — No Charge", 32, TEAL, True, font_name="Outfit")

card(slide, 0.8, 2.3, 5.5, 3.8)
txt(slide, 1.1, 2.5, 5, 0.4, "immersive-family-dentistry.netlify.app", 15, LJ_BLUE, True)
bullets(slide, 1.1, 3.0, 5, 2.8, [
    "Full responsive website with your branding",
    "Team profiles, technology showcase",
    "Learning Center with your publications",
    "",
    "This is a gift — no strings attached.",
    "I respect your existing web contract."], 13, DARK_TEXT)

card(slide, 7.0, 2.3, 5.5, 3.8)
txt(slide, 7.3, 2.5, 5, 0.4, "But the bigger idea...", 15, EMERALD, True)
bullets(slide, 7.3, 3.0, 5, 2.8, [
    "The website showcases the vision.",
    "The software IS the business.",
    "",
    "We call it MouthMap —",
    "and it is what I really want",
    "to talk about today."], 13, TEAL)

add_logo(slide)

notes(slide, "I built this website over the weekend. It is yours, free. I know you have a contract — I respect that. But it became the canvas for the bigger idea: MouthMap. The website is the showcase. The software is the business. Let me show you the problem it solves.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: THE PROBLEM (1.5 min)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_SECTION)

txt(slide, 0.8, 1.2, 5, 0.3, "THE PROBLEM", 14, GOLD, True)
txt(slide, 0.8, 1.55, 10, 0.5, "Patients Don't Understand Their Oral Health", 32, TEAL, True, font_name="Outfit")

card(slide, 0.8, 2.3, 5.5, 3.8)
txt(slide, 1.1, 2.5, 5, 0.3, "The Reality Today", 17, TEAL, True)
bullets(slide, 1.1, 2.9, 5, 3.0, [
    "Average case acceptance: only 34%",
    "Patients forget 80% of chairside discussion",
    "Scans and AI data stay locked in PMS",
    "Family decision-makers never in the chair",
    "No way to review findings at home",
    "Result: lost revenue + poorer outcomes"], 13, DARK_TEXT)

card(slide, 7.0, 2.3, 5.5, 3.8)
txt(slide, 7.3, 2.5, 5, 0.3, "The Opportunity", 17, EMERALD, True)
bullets(slide, 7.3, 2.9, 5, 3.0, [
    "You already capture 3D scans (Medit)",
    "AI already analyzes X-rays in real-time",
    "The data exists — patients can't access it",
    "Practices with visual tools: 60-85% acceptance",
    "That's a 26-51% improvement",
    "Understanding = trust = case acceptance"], 13, DARK_TEXT)

notes(slide, "34% case acceptance. 80% forgotten within a week. PMS stands for Practice Management Software — Dentrix, Eaglesoft. The data stays locked there. But practices that use visual presentation tools see 60-85% acceptance. The data exists in your workflow. The gap is getting it to the patient at home.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: WHAT IS MOUTHMAP (2 min)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_SECTION)

txt(slide, 0.8, 1.2, 5, 0.3, "THE PRODUCT", 14, GOLD, True)
txt(slide, 0.8, 1.55, 10, 0.5, "MouthMap — The Open-Architecture Patient Portal", 32, TEAL, True, font_name="Outfit")
txt(slide, 0.8, 2.1, 11, 0.3, "A HIPAA-compliant portal for ANY scanner ecosystem. Patient scans + AI findings + messaging — from any device.", 14, MUTED_TEXT)

features = [
    ("\U0001F9B7", "3D Scan Explorer", "Patients rotate and explore\ntheir own 3D mouth scan.\nTap teeth to see findings."),
    ("\U0001F4CA", "AI Health Report", "Plain-language summary.\nGreen (healthy), yellow\n(monitor), red (treat)."),
    ("\U0001F4C8", "Progress Timeline", "Before-and-after 3D\ncomparisons across visits.\nVisual proof of progress."),
    ("\U0001F4AC", "Care Hub", "Secure messaging.\nTreatment reminders.\nEducational content.")
]

for i, (icon, title, desc) in enumerate(features):
    x = 0.5 + i * 3.15
    card(slide, x, 2.6, 2.95, 3.5)
    txt(slide, x + 0.2, 2.7, 2.5, 0.4, icon, 28, DARK_TEXT, False, PP_ALIGN.CENTER)
    txt(slide, x + 0.2, 3.15, 2.5, 0.3, title, 14, TEAL, True, PP_ALIGN.CENTER)
    txt(slide, x + 0.2, 3.55, 2.5, 2.2, desc, 11, MUTED_TEXT, False, PP_ALIGN.CENTER)

add_logo(slide, 11.0, 2.7, 0.8)

notes(slide, """HIPAA stands for Health Insurance Portability and Accountability Act. MouthMap is a secure web portal — no app download.

How it works: After a cleaning Monday, patient gets an email Tuesday — 'Your MouthMap is ready.' They log in, explore their 3D scan, read the AI summary in plain English, and share with their spouse. Data flows automatically from Medit Link and AI platforms through APIs — Application Programming Interfaces. Zero manual uploads. Zero extra clicks for your team.

By their next visit, they have already made the treatment decision.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: WHY NOT JUST BUY 3SHAPE? (2 min)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_SECTION)

txt(slide, 0.8, 1.2, 5, 0.3, "THE HONEST QUESTION", 14, GOLD, True)
txt(slide, 0.8, 1.55, 10, 0.5, "Why Not Just Buy 3Shape DentalHealth?", 32, TEAL, True, font_name="Outfit")
txt(slide, 0.8, 2.1, 11, 0.3, "3Shape proves the concept works. A major company invested millions. That is market validation.", 14, MUTED_TEXT)

card(slide, 0.8, 2.5, 5.5, 3.8)
txt(slide, 1.1, 2.7, 5, 0.3, "3Shape DentalHealth (Closed)", 15, MUTED_TEXT, True)
bullets(slide, 1.1, 3.1, 5, 3.0, [
    "Requires TRIOS 6 scanner ($20K+)",
    "You would abandon Medit + KOL status",
    "Only 3Shape's AI — no Pearl, Overjet,",
    "  or VideaHealth integration",
    "Cannot customize or rebrand",
    "Cannot license to other practices",
    "You are a customer, not an owner",
    "No IP, no equity, no business"], 12, DARK_TEXT)

card(slide, 7.0, 2.5, 5.5, 3.8)
txt(slide, 7.3, 2.7, 5, 0.3, "MouthMap (Open Architecture)", 15, TEAL, True)
bullets(slide, 7.3, 3.1, 5, 3.0, [
    "Works with Medit, iTero, ANY scanner",
    "Preserves your KOL relationships",
    "Integrates any FDA-cleared AI platform",
    "  (Pearl, Overjet, VideaHealth, etc.)",
    "Fully customizable per practice",
    "White-label to DSOs & manufacturers",
    "You are a co-owner of the IP",
    "Equity + revenue streams"], 12, DARK_TEXT)

notes(slide, """I want to be upfront. 3Shape already built something similar — their DentalHealth App for TRIOS 6 users. It is good and it proves the market is real.

But three problems for you. First, hardware lock-in — you would have to buy a TRIOS 6 for twenty thousand plus and abandon Medit. That means walking away from your KOL — Key Opinion Leader — relationships with Medit and SprintRay. That is a career cost, not just a financial one.

Second, closed ecosystem. 3Shape only uses 3Shape's AI. MouthMap integrates any FDA — Food and Drug Administration — cleared AI platform.

Third — with 3Shape, you are a customer. With MouthMap, you are a co-owner of the intellectual property. You can license it to DSOs — Dental Service Organizations — like Aspen and Heartland, or white-label it.

The analogy: 3Shape is the iPhone, locked to their hardware. MouthMap is Android — open platform, any hardware.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: ZERO EXTRA WORK (1.5 min)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_SECTION)

txt(slide, 0.8, 1.2, 5, 0.3, "YOUR CONCERN, ADDRESSED", 14, GOLD, True)
txt(slide, 0.8, 1.55, 10, 0.5, "\"Dentists Don't Want Extra Work\"", 32, TEAL, True, font_name="Outfit")
txt(slide, 0.8, 2.1, 10, 0.3, "You said it — and you're right. MouthMap creates ZERO extra work for the provider.", 14, MUTED_TEXT)

teal_card(slide, 0.8, 2.5, 5.5, 3.5)
txt(slide, 1.1, 2.7, 5, 0.3, "Your Workflow (Unchanged)", 16, GOLD, True)
bullets(slide, 1.1, 3.1, 5, 2.6, [
    "1. Scan with Medit i700 → Medit Link",
    "2. AI runs on X-rays → auto-generated",
    "3. Treatment plan created in PMS",
    "4. Patient checks out and leaves",
    "",
    "Nothing changes for you or your team."], 12, WHITE)

card(slide, 7.0, 2.5, 5.5, 3.5)
txt(slide, 7.3, 2.7, 5, 0.3, "What MouthMap Does Automatically", 16, EMERALD, True)
bullets(slide, 7.3, 3.1, 5, 2.6, [
    "1. Medit scan auto-syncs to portal",
    "2. AI findings populate health summary",
    "3. Patient gets: \"Your MouthMap is ready\"",
    "4. Patient reviews at home with family",
    "",
    "Zero clicks. Zero uploads. Zero extra time."], 12, DARK_TEXT)

notes(slide, """This addresses your exact concern from your email. You said, quote, 'Most dentists are busy and they do not like the extra work.'

You are right. So MouthMap pulls data automatically through APIs — Application Programming Interfaces — from Medit Link and your AI platform. Your workflow stays identical. The system works because it pulls data from tools you already use — it does not ask you to push data somewhere new.

Also: estimated 45-60 minutes saved per day from fewer re-explanations and reduced phone calls. That is recoverable chair time.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: REVENUE MODEL (1.5 min)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_SECTION)

txt(slide, 0.8, 1.2, 5, 0.3, "THE BUSINESS", 14, GOLD, True)
txt(slide, 0.8, 1.55, 10, 0.5, "Three Ways MouthMap Generates Revenue", 32, TEAL, True, font_name="Outfit")

# Revenue stream 1
card(slide, 0.5, 2.3, 4.0, 3.8)
txt(slide, 0.8, 2.5, 3.4, 0.3, "1. Your Practice First", 16, TEAL, True)
txt(slide, 0.8, 2.9, 3.4, 0.25, "Prove it works at Lake Jeanette", 11, GOLD, True)
bullets(slide, 0.8, 3.3, 3.4, 2.6, [
    "Increase case acceptance",
    "Reduce re-explanations",
    "Improve patient retention",
    "Measurable ROI",
    "",
    "If it works here,",
    "it works anywhere."], 11, DARK_TEXT)

# Revenue stream 2
card(slide, 4.8, 2.3, 4.0, 3.8)
txt(slide, 5.1, 2.5, 3.4, 0.3, "2. SaaS Licensing", 16, TEAL, True)
txt(slide, 5.1, 2.9, 3.4, 0.25, "License to other practices", 11, GOLD, True)
bullets(slide, 5.1, 3.3, 3.4, 2.6, [
    "$500-$1,500/mo per practice",
    "200,000+ U.S. dental practices",
    "Your KOL network = distribution",
    "50 practices = $300K-$900K ARR",
    "",
    "SaaS = Software as a Service.",
    "Valued at 10x annual revenue."], 11, DARK_TEXT)

# Revenue stream 3
card(slide, 9.1, 2.3, 4.0, 3.8)
txt(slide, 9.4, 2.5, 3.4, 0.3, "3. White-Label IP", 16, TEAL, True)
txt(slide, 9.4, 2.9, 3.4, 0.25, "License to DSOs & manufacturers", 11, GOLD, True)
bullets(slide, 9.4, 3.3, 3.4, 2.6, [
    "DSOs (Aspen, Heartland)",
    "  rebrand as their own portal",
    "Medit could bundle it",
    "  with their scanners",
    "Insurance co's as benefit",
    "",
    "Highest margin revenue."], 11, DARK_TEXT)

notes(slide, """Three revenue streams. First — deploy at Lake Jeanette, prove the ROI — return on investment — with real case acceptance data.

Second — SaaS — Software as a Service — licensing to other practices at $500-$1,500/month. Your KOL — Key Opinion Leader — network is the distribution channel. ARR stands for Annual Recurring Revenue. SaaS companies are valued at 10x revenue.

Third — white-label. DSOs — Dental Service Organizations — like Aspen and Heartland rebrand it as their own. Scanner manufacturers bundle it with hardware. We own the intellectual property and license it.

Funding: NIH SBIR — National Institutes of Health Small Business Innovation Research — grants cover development. $275K Phase 1, up to $1M Phase 2. Non-dilutive — no equity given up. Think! Ventures Foundation, my 501(c)(3), applies for the grants. Zero cost to you.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: WHAT I NEED FROM YOU (Transition to conversation)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_SECTION)

txt(slide, 0.8, 1.2, 5, 0.3, "OVER TO YOU", 14, GOLD, True)
txt(slide, 0.8, 1.55, 10, 0.5, "Your Feedback Shapes Everything", 32, TEAL, True, font_name="Outfit")
txt(slide, 0.8, 2.1, 10, 0.3, "I need your honest, in-depth feedback to brief my developer and ground the grant application in clinical reality.", 14, MUTED_TEXT)

card(slide, 0.8, 2.5, 11.5, 3.6)

questions = [
    "1. Walk me through your digital workflow: Scan → AI → PMS → what happens to the data?",
    "2. What data do you wish patients could see from home?",
    "3. What's the biggest gap between what you capture and what patients understand?",
    "4. SmileCloud — complement or compete? How do you use it today?",
    "5. What would make you willing to pilot this with patients?",
    "6. What's the #1 objection other dentists would raise?",
    "7. Existing research at HPU — could MouthMap fit alongside?",
    "8. Could HPU students participate in beta testing or an IRB study?",
]

bullets(slide, 1.3, 2.7, 10.5, 3.2, questions, 14, DARK_TEXT)

notes(slide, """This is where I stop presenting and start listening. These questions shape what I tell Gary Glass — my software developer — to build.

HPU stands for High Point University. IRB stands for Institutional Review Board — the ethics committee that approves research with human subjects.

Be candid. If this will not work, tell me why. That is just as valuable.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: NEXT STEPS (30 sec — reference if conversation goes well)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_SECTION)

txt(slide, 0.8, 1.2, 5, 0.3, "IF THIS RESONATES", 14, GOLD, True)
txt(slide, 0.8, 1.55, 10, 0.5, "Next Steps — Every Step Costs $0", 32, TEAL, True, font_name="Outfit")

steps_data = [
    ("Today", "Your feedback\non this concept", "$0"),
    ("June", "Shadow your workflow,\ndocument data flow", "$0"),
    ("July", "Technical feasibility\nreport (API access?)", "$0"),
    ("Aug-Sep", "NIH SBIR grant\napplication filed", "$0"),
    ("Q1 2027", "Beta launch at\nLake Jeanette", "$0 (grant)")
]

for i, (when, what, cost) in enumerate(steps_data):
    x = 0.3 + i * 2.55
    card(slide, x, 2.3, 2.35, 3.5)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.75), Inches(2.45), Inches(0.7), Inches(0.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = TEAL; circle.line.fill.background()
    txt(slide, x + 0.1, 2.5, 2.15, 0.55, when, 16, WHITE, True, PP_ALIGN.CENTER)
    txt(slide, x + 0.1, 3.3, 2.15, 1.5, what, 12, DARK_TEXT, False, PP_ALIGN.CENTER)
    txt(slide, x + 0.1, 4.8, 2.15, 0.5, cost, 20, EMERALD, True, PP_ALIGN.CENTER)

notes(slide, """Only show this if the conversation is going well. Every step costs zero.

NIH SBIR = National Institutes of Health Small Business Innovation Research grant. Think! Ventures Foundation — my 501(c)(3) nonprofit — applies. You provide a letter of support. MVP = Minimum Viable Product — the simplest working version we test with real patients.

The team: Me (product/strategy), Gary Glass (software engineering), and you (clinical guidance + flagship practice). Three people, lean by design.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: CLOSING (Leave up during Q&A)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_SECTION)

txt(slide, 1.5, 1.5, 10, 0.5, "THANK YOU", 16, GOLD, True, PP_ALIGN.CENTER)
txt(slide, 1.5, 2.1, 10, 1.0,
    "\"3Shape proved the concept.\nThe open-architecture version\ndoesn't exist yet.\"",
    28, TEAL, False, PP_ALIGN.CENTER)

txt(slide, 1.5, 3.6, 10, 0.5, "MouthMap — Navigate Your Oral Health.", 24, GOLD, True, PP_ALIGN.CENTER)
txt(slide, 1.5, 4.3, 10, 0.4, "W. Chris Harrison  |  Think! Design and Planning, LLC", 16, MUTED_TEXT, False, PP_ALIGN.CENTER)
txt(slide, 1.5, 4.8, 10, 0.4, "immersive-family-dentistry.netlify.app", 14, LJ_BLUE, False, PP_ALIGN.CENTER)

if os.path.exists(MM_LOGO):
    slide.shapes.add_picture(MM_LOGO, Inches(5.8), Inches(5.0), height=Inches(1.3))

notes(slide, """Leave this slide up during open conversation. The closing quote reframes the entire pitch: 3Shape validated the concept. The open-architecture version for the rest of the industry does not exist yet. That is the opportunity.

Thank you, Dr. Brenes.""")

# ═══════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
print("V7 — TIGHT 10 slides. ~10 min present, 35 min conversation.")
