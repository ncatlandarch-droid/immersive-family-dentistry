from pptx import Presentation
prs = Presentation(r'C:\Users\Chris\Desktop\WEBSITES\THINK! VENTURES\docs\Think_Ventures_Template_converted.pptx')
layout = prs.slide_layouts[2]
print('=== Section Header Layout - ALL shapes ===')
for shape in layout.shapes:
    print(f'  {shape.name}: type={shape.shape_type}')
    print(f'    pos (emu): left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}')
    print(f'    pos (in): left={shape.left/914400:.2f}, top={shape.top/914400:.2f}, w={shape.width/914400:.2f}, h={shape.height/914400:.2f}')
    if hasattr(shape, 'text') and shape.text:
        print(f'    text: "{shape.text}"')
        if hasattr(shape, 'text_frame'):
            for p in shape.text_frame.paragraphs:
                f = p.font
                try:
                    c = f.color.rgb if f.color and f.color.rgb else "inherit"
                except:
                    c = "inherit"
                print(f'    font: name={f.name}, size={f.size}, bold={f.bold}, color={c}')
    if shape.shape_type == 9:
        try:
            print(f'    line color: {shape.line.color.rgb}')
        except:
            print(f'    line color: inherit')
        print(f'    line width: {shape.line.width}')
print(f'\nSlide dims: {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f} in')
